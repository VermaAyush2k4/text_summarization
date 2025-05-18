from flask import Flask, render_template, request, jsonify
from transformers import pipeline
import os
from werkzeug.utils import secure_filename
import PyPDF2
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation

app = Flask(__name__)
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

# Configure upload folder
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'txt', 'doc', 'docx', 'pdf', 'ppt', 'pptx'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    try:
        # First try PyPDF2
        pdf_reader = PyPDF2.PdfReader(file_path)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text
    except Exception:
        # If PyPDF2 fails, try pdfminer
        try:
            return extract_text(file_path)
        except Exception as e:
            raise Exception(f"Error extracting text from PDF: {str(e)}")

def extract_text_from_docx(file_path):
    """Extract text from DOCX file"""
    try:
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        raise Exception(f"Error extracting text from DOCX: {str(e)}")

def extract_text_from_pptx(file_path):
    """Extract text from PPTX file"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        raise Exception(f"Error extracting text from PPTX: {str(e)}")

def extract_text_from_file(file_path):
    """Extract text from any supported file type"""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.doc', '.docx']:
        return extract_text_from_docx(file_path)
    elif ext in ['.ppt', '.pptx']:
        return extract_text_from_pptx(file_path)
    elif ext == '.txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    else:
        raise Exception(f"Unsupported file type: {ext}")

def summarize_text(text, max_chunk=1024, slider_position=5):
    """
    Summarize text using BART model
    Args:
        text: Input text
        max_chunk: Maximum chunk size for processing
        slider_position: Slider position (1-10) that determines summary length
    """
    text = text.strip().replace("\n", " ")
    summaries = []
    
    # Calculate base summary length based on slider position
    input_words = len(text.split())
    base_summary_length = input_words // 3  # 1/3 of input length
    
    # Apply slider position to get final summary length
    summary_length = int((slider_position / 10) * base_summary_length)
    
    # Ensure minimum length of 40 words
    summary_length = max(40, summary_length)
    
    # Ensure maximum length of 1000 words
    summary_length = min(1000, summary_length)
    
    min_length = max(40, int(summary_length * 0.7))
    max_length = min(1000, int(summary_length * 1.3))
    
    while len(text) > max_chunk:
        split_at = text[:max_chunk].rfind(".")
        if split_at == -1:
            split_at = max_chunk
        chunk = text[:split_at+1]
        text = text[split_at+1:]
        summary = summarizer(chunk, 
                           max_length=max_length, 
                           min_length=min_length, 
                           do_sample=False)[0]['summary_text']
        summaries.append(summary)
    
    if text:
        summary = summarizer(text, 
                           max_length=max_length, 
                           min_length=min_length, 
                           do_sample=False)[0]['summary_text']
        summaries.append(summary)
    
    return " ".join(summaries)

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/summarize", methods=["POST"])
def summarize():
    data = request.get_json()
    text = data.get("text", "")
    slider_position = int(data.get("slider_position", 5))
    
    if not text.strip():
        return jsonify({"summary": "Please enter some text to summarize."})
    
    try:
        summary = summarize_text(text, slider_position=slider_position)
        return jsonify({"summary": summary})
    except Exception as e:
        return jsonify({"summary": "An error occurred during summarization.", "error": str(e)}), 500

@app.route("/upload", methods=["POST"])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        try:
            # Extract text based on file type
            content = extract_text_from_file(file_path)
            return jsonify({"text": content})
        except Exception as e:
            return jsonify({"error": f"Error processing file: {str(e)}"}), 500
    
    return jsonify({"error": "File type not allowed"}), 400

if __name__ == "__main__":
    app.run(debug=True)
